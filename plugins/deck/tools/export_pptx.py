"""deck.export_pptx：真实导出——deck.presentation workflow 的 export 段 provider。

读最新页面文档 + 视觉规范 → python-pptx 真实构建 .pptx（封面/内容/结尾，要点与演讲
备注齐全，底色/强调色按视觉规范）→ 重新打开校验页数 → 落盘 exports/<deck_id>/v<N>.pptx
+ deck_exports 表 + artifact deck.export.pptx（ref=deck_id）+ rendered_from 边 →
deck.document（设计文档 §4.4 合法关系）。

「可打开导出」由重开 python-pptx 读回校验（slide 数 == 文档页数），不是工具自报成功。
文件自包含（加载器按文件独立 importlib 加载，禁止跨文件 import）。
"""
import json
import os
import time
from pathlib import Path

from yibao_brain.ipc import ActionResult, RiskLevel
from yibao_brain.tools import Tool

_SLIDE_W = 12192000   # 16:9 宽屏（EMU：13.333in）
_SLIDE_H = 6858000    # 7.5in


def _hex(color: str) -> tuple[int, int, int]:
    c = color.lstrip("#")
    return int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)


class ExportPptx(Tool):
    id = "deck.export_pptx"
    label = "导出 pptx"
    description = (
        "把最新页面文档导出成真实 .pptx：封面 + 内容页（标题 + 要点 + 演讲备注）+ 结尾，"
        "底色/强调色按视觉规范。导出后重新打开校验页数，不一致即失败。"
        "落盘 exports/<deck_id>/v<N>.pptx（版本递增）。"
    )
    default_risk = RiskLevel.L2_MEDIUM
    work_outputs = (
        {
            "kind": "artifact",
            "artifact_type": "deck.export.pptx",
            "ref_from": "data.export_ref",
            "content_ref_from": "data.path",
            "metadata_fields": ["data.version", "data.slide_count", "data.path"],
        },
        {   # 导出 rendered_from 页面文档（导出来源血缘）
            "kind": "edge",
            "relation": "rendered_from",
            "source_artifact_type": "deck.export.pptx",
            "source_ref_from": "data.export_ref",
            "target_artifact_type": "deck.document",
            "target_ref_from": "data.deck_ref",
        },
    )

    def __init__(self, data_dir: str):
        self._plugin_root = Path(data_dir)

    def openai_schema(self) -> dict:
        return {
            "name": self.id,
            "description": self.description,
            "parameters": {
                "type": "object",
                "properties": {
                    "deck_id": {"type": "string", "description": "演示文稿 id"},
                },
                "required": ["deck_id"],
            },
        }

    def run(self, params: dict, ctx) -> ActionResult:
        deck_id = str(params.get("deck_id") or "").strip()
        if not deck_id:
            return ActionResult(success=False, error="deck_id 不能为空")
        if not ctx.db.query("decks", where={"id": deck_id}):
            return ActionResult(success=False, error=f"演示文稿不存在：{deck_id}")
        doc, error = _load_latest(ctx, deck_id, "document")
        if error:
            return ActionResult(success=False, error=error)
        visual, _ = _load_latest(ctx, deck_id, "visual")  # 没有视觉规范用缺省
        try:
            from pptx import Presentation  # noqa: F401 —— 依赖检查前置
        except ImportError:
            return ActionResult(success=False, error="缺少 python-pptx（sidecar 依赖）")
        latest = ctx.db.query("deck_exports", where={"deck_id": deck_id},
                              order="version DESC", limit=1)
        version = (int(latest[0]["version"]) if latest else 0) + 1
        out_dir = self._plugin_root / "exports" / deck_id
        try:
            out_dir.mkdir(parents=True, exist_ok=True)
        except OSError as e:
            return ActionResult(success=False, error=f"创建导出目录失败：{e}")
        out_path = out_dir / f"v{version}.pptx"
        error = _build_pptx(doc, visual, out_path)
        if error:
            return ActionResult(success=False, error=error)
        # 验收：重开文件读回页数（可打开导出，不看工具自报）
        slide_count, error = _verify_pptx(out_path, len(doc.get("slides") or []))
        if error:
            return ActionResult(success=False, error=error)
        now = int(time.time())
        ctx.db.insert("deck_exports", {
            "deck_id": deck_id, "version": version, "path": str(out_path),
            "slide_count": slide_count, "created_at": now,
        })
        return ActionResult(success=True, data={
            "deck_id": deck_id,
            "deck_ref": deck_id,
            "export_ref": deck_id,
            "version": version,
            "path": str(out_path),
            "slide_count": slide_count,
        })


def _load_latest(ctx, deck_id: str, kind: str) -> tuple[dict, str]:
    rows = ctx.db.query("deck_stages", where={"deck_id": deck_id, "kind": kind},
                        order="version DESC", limit=1)
    if not rows:
        return {}, ""
    blobs = getattr(ctx, "blobs", None)
    if blobs is None:
        return {}, "底座未提供 blobs capability"
    try:
        path = blobs.resolve(str(rows[0]["content_path"]), require_exists=False)
        doc = json.loads(Path(path).read_text(encoding="utf-8"))
    except (OSError, ValueError) as e:
        return {}, f"{kind} 读取失败：{e}"
    return doc, ""


def _build_pptx(doc: dict, visual: dict, out_path: Path) -> str:
    """构建真实 .pptx。返回 error（空串=成功）。"""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Emu, Pt
    except ImportError:
        return "缺少 python-pptx"
    palette = (visual or {}).get("palette") or "#111418"
    accent = (visual or {}).get("accent") or "#5EA0D2"
    bg = RGBColor(*_hex(palette))
    fg = RGBColor(245, 245, 240)
    ac = RGBColor(*_hex(accent))
    try:
        prs = Presentation()
        prs.slide_width = Emu(_SLIDE_W)
        prs.slide_height = Emu(_SLIDE_H)
        blank = prs.slide_layouts[6]  # 空白版式，全自绘
        for slide_spec in doc.get("slides") or []:
            slide = prs.slides.add_slide(blank)
            # 底色
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = bg
            # 标题
            title_text = str(slide_spec.get("title") or "").strip() or "（无标题）"
            box = slide.shapes.add_textbox(Emu(609600), Emu(457200), Emu(_SLIDE_W - 1219200), Emu(1371600))
            tf = box.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.text = title_text
            para.font.size = Pt(40 if slide_spec.get("kind") == "title" else 32)
            para.font.bold = True
            para.font.color.rgb = ac if slide_spec.get("kind") == "title" else fg
            # 内容要点
            bullets = [str(b) for b in (slide_spec.get("bullets") or [])]
            if bullets:
                body = slide.shapes.add_textbox(Emu(609600), Emu(2057400), Emu(_SLIDE_W - 1219200), Emu(4114800))
                btf = body.text_frame
                btf.word_wrap = True
                for i, bullet in enumerate(bullets):
                    p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                    p.text = "· " + bullet
                    p.font.size = Pt(20)
                    p.font.color.rgb = fg
                    p.space_after = Pt(10)
            # 演讲备注
            note = str(slide_spec.get("note") or "").strip()
            if note:
                slide.notes_slide.notes_text_frame.text = note
        prs.save(str(out_path))
        return ""
    except (OSError, ValueError, RuntimeError) as e:
        return f"pptx 构建失败：{e}"


def _verify_pptx(out_path: Path, expected_slides: int) -> tuple[int, str]:
    """重开导出文件校验页数。返回 (slide_count, error)。"""
    try:
        from pptx import Presentation

        count = len(Presentation(str(out_path)).slides)  # 重开读回：可打开导出才算数
        if count != expected_slides:
            return count, f"导出校验失败：读回 {count} 页 ≠ 文档 {expected_slides} 页"
        return count, ""
    except Exception as e:
        return 0, f"导出文件不可读：{e}"


def make_tools(ctx):
    return [ExportPptx(os.path.dirname(ctx.db.path))]
