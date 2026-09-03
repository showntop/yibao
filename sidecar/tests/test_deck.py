"""deck 插件（演示文稿链）：内核泛化验证——视频之外第二个领域用同一内核跑通。

核心断言：
- 全链：deck_save → claims_save → storyline_save → compose → visual_save → validate
  → export_pptx，每段都落真实 artifact（blob/文件）+ Work Graph 事件；
- deck.presentation 预检：deck 插件加载后八段全 available（内核泛化的验收锚点）；
- validate 诚实卡门：占位符残留/要点超限时 success=False；
- export_pptx 产出真实 .pptx（python-pptx 重开读回页数 == 文档页数）；
- 失败前置：缺 deck_id / 缺上游产物的人话报错。
"""
import json
from pathlib import Path

import pytest

from yibao_brain.durable_execution import DurableExecutionEngine
from yibao_brain.llm import FakeProvider
from yibao_brain.memory import FakeMemory
from yibao_brain.plugins import LlmChat, load_plugins
from yibao_brain.tools import ToolRegistry
from yibao_brain.work_graph import WorkGraphStore, build_capability_index

REPO_ROOT = Path(__file__).resolve().parents[2]

needs_pptx = pytest.mark.skipif(
    __import__("importlib.util").util.find_spec("pptx") is None,
    reason="需要 python-pptx（sidecar 依赖）",
)


@pytest.fixture
def data_dir(tmp_path, monkeypatch):
    d = tmp_path / "data"
    monkeypatch.setenv("YIBAO_DATA_DIR", str(d))
    return d


@pytest.fixture
def env(data_dir, tmp_path):
    reg = ToolRegistry()

    class _Http:
        def get(self, url, **kw):
            return {}

        def post(self, url, **kw):
            return {}

    load_plugins(
        REPO_ROOT / "plugins", reg,
        memory=FakeMemory(), http=_Http(), llm=LlmChat(FakeProvider()),
        durable_engine=DurableExecutionEngine(WorkGraphStore(str(tmp_path / "wg.db"))),
    )
    return reg


def _run(reg, tid, params):
    t = reg.get(tid)
    return t.run(params, t.plugin_ctx)


def _make_deck(reg, title="Q3 策略汇报") -> str:
    r = _run(reg, "deck.deck_save", {
        "title": title, "audience": "管理层", "goal": "过 Q3 策略", "page_count": 6,
    })
    assert r.success, r.error
    return r.data["deck_id"]


def _drive_to_document(reg, deck_id):
    assert _run(reg, "deck.claims_save", {
        "deck_id": deck_id,
        "claims": [
            {"claim": "增长来自留存", "support": "Q2 留存环比 +18%", "source_uri": ""},
            {"claim": "渠道 B 贡献过半", "support": "渠道 B 占新增 54%", "source_uri": ""},
        ],
    }).success
    assert _run(reg, "deck.storyline_save", {
        "deck_id": deck_id,
        "sections": [
            {"title": "现状", "key_points": ["Q2 留存 +18%", "渠道 B 占新增 54%"]},
            {"title": "问题", "key_points": ["拉新成本翻倍", "老客复购下滑"]},
            {"title": "策略", "key_points": ["砍低效投放", "加码渠道 B", "会员体系"]},
        ],
    }).success
    r = _run(reg, "deck.compose", {"deck_id": deck_id})
    assert r.success, r.error
    return r


# ---------- 立项与预检（内核泛化锚点） ----------


def test_deck_plugin_covers_deck_presentation_chain(env, tmp_path):
    """deck.presentation 预检：deck 插件加载后八段全 available——同一内核，
    不改 OS schema，第二个领域（演示文稿）跑通。这是内核泛化的验收锚点。"""
    reg = env
    index = build_capability_index(reg.list())
    for artifact_type in ("deck.brief", "deck.claim_set", "deck.storyline", "deck.document",
                          "deck.visual_spec", "quality.report", "deck.export.pptx"):
        assert artifact_type in index, artifact_type

    graph = WorkGraphStore(str(tmp_path / "work_graph.db"))
    try:
        graph.set_capability_providers(index)
        graph.create_workspace("ws", "Q3 策略汇报 pptx", str(tmp_path / "ws"))
        run = graph.workspace_view("ws")["workflow_run"]
        assert run["definition_id"] == "deck.presentation"
        assert run["status"] == "ready"  # 八段全可满足：不 blocked
        assert run["capability_plan"]["missing"] == []
    finally:
        graph.close()


# ---------- 全链 ----------


@needs_pptx
def test_deck_full_chain_to_real_pptx(env, data_dir):
    """立项 → 主张 → 故事线 → 组装 → 视觉 → 校验 → 导出：真实 .pptx 落盘且读回页数一致。"""
    reg = env
    deck_id = _make_deck(reg)
    _drive_to_document(reg, deck_id)
    assert _run(reg, "deck.visual_save", {"deck_id": deck_id, "palette": "#0E2A47"}).success
    rv = _run(reg, "deck.validate", {"deck_id": deck_id})
    assert rv.success and rv.data["ok"] is True
    rx = _run(reg, "deck.export_pptx", {"deck_id": deck_id})
    assert rx.success, rx.error
    path = Path(rx.data["path"])
    assert path.is_file() and path.suffix == ".pptx" and data_dir in path.parents
    assert rx.data["slide_count"] == 5  # 封面 + 3 内容页 + 结尾
    # 读回内容真实可验：标题页带演示标题
    from pptx import Presentation

    prs = Presentation(str(path))
    texts = [shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame]
    assert any("Q3 策略汇报" in t for t in texts)
    # 演讲备注随页落盘
    content_slide = prs.slides[1]
    assert "本段主张" in content_slide.notes_slide.notes_text_frame.text


@needs_pptx
def test_deck_validate_fails_honestly_on_placeholder(env):
    """校验卡门诚实：文档里留「待补」占位符 → success=False + 报告明细，不蒙混。"""
    reg = env
    deck_id = _make_deck(reg)
    assert _run(reg, "deck.claims_save", {"deck_id": deck_id, "claims": [
        {"claim": "增长来自留存", "support": "留存 +18%"},
    ]}).success
    assert _run(reg, "deck.storyline_save", {"deck_id": deck_id, "sections": [
        {"title": "现状", "key_points": ["留存 +18%", "待补"]},
    ]}).success
    assert _run(reg, "deck.compose", {"deck_id": deck_id}).success
    r = _run(reg, "deck.validate", {"deck_id": deck_id})
    assert not r.success
    assert "占位符" in r.error
    # 校验不过也落了报告 artifact（failure 也是事实）
    assert r.data["ok"] is False
    failed = [c for c in r.data["checks"] if not c["ok"]]
    assert failed and "占位符" in failed[0]["name"]


def test_deck_missing_prerequisites_error_in_human_terms(env):
    """缺上游产物/缺 deck_id：人话报错，不抛栈。"""
    reg = env
    assert not _run(reg, "deck.compose", {"deck_id": "ghost"}).success
    deck_id = _make_deck(reg)
    r = _run(reg, "deck.compose", {"deck_id": deck_id})  # 没存故事线
    assert not r.success and "故事线" in r.error
    r = _run(reg, "deck.validate", {"deck_id": deck_id})  # 没组装
    assert not r.success and "组装" in r.error
    assert not _run(reg, "deck.deck_save", {"title": ""}).success


@needs_pptx
def test_deck_workflow_run_reaches_completed_via_artifacts(env, tmp_path):
    """验收驱动的关键证据：产物经 invoker + WorkGraphInvocationSink 推进
    deck.presentation 到 completed——没有 export artifact 时 run 不 completed。

    与视频链（test_zimeiti_video_chain）同一条投影路径：换领域不换内核。
    """
    from yibao_brain.audit import AuditLog
    from yibao_brain.invoker import ToolInvoker
    from yibao_brain.safety import Gate, GatePolicy, RiskClassifier
    from yibao_brain.llm import ToolCall
    from yibao_brain.work_events import WorkGraphInvocationSink

    reg = env
    graph = WorkGraphStore(str(tmp_path / "work_graph.db"))
    try:
        graph.set_capability_providers(build_capability_index(reg.list()))
        graph.create_workspace("ws", "Q3 策略汇报 pptx", str(tmp_path / "ws"))
        run = graph.workspace_view("ws")["workflow_run"]
        assert run["status"] == "ready"

        invoker = ToolInvoker(
            reg, RiskClassifier(), Gate(GatePolicy()), AuditLog(str(tmp_path / "audit.db")),
        )
        invoker.invocation_sink = WorkGraphInvocationSink(graph, lambda _cid: "ws")

        def execute(tool_id, params):
            action = invoker.propose(ToolCall(id=f"tc-{tool_id}", tool_id=tool_id, params=params))
            result = invoker.execute(action, params, {"conversation_id": "s1", "surface": "home"})
            return result

        deck_id = execute("deck.deck_save", {
            "title": "Q3 策略汇报", "audience": "管理层", "goal": "过 Q3 策略", "page_count": 6,
        }).data["deck_id"]
        execute("deck.claims_save", {"deck_id": deck_id, "claims": [
            {"claim": "增长来自留存", "support": "留存 +18%"}]})
        execute("deck.storyline_save", {"deck_id": deck_id, "sections": [
            {"title": "现状", "key_points": ["留存 +18%"]}]})
        execute("deck.compose", {"deck_id": deck_id})
        execute("deck.visual_save", {"deck_id": deck_id})
        execute("deck.validate", {"deck_id": deck_id})

        # 还差 export：run 不得 completed
        run = graph.workspace_view("ws")["workflow_run"]
        assert run["status"] != "completed"

        r = execute("deck.export_pptx", {"deck_id": deck_id})
        assert r.success, r.error
        run = graph.workspace_view("ws")["workflow_run"]
        assert run["status"] == "completed"
    finally:
        graph.close()
